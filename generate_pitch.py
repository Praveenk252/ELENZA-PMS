#!/usr/bin/env python3
"""
Elenza PMS — 100-Slide Pitch Deck Generator
Run:  python generate_pitch.py
Output: ElenzaPMS_PitchDeck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

# ─── Brand Palette ────────────────────────────────────────────────
BLUE       = RGBColor(0x04, 0x6B, 0xD2)
BLUE_DARK  = RGBColor(0x04, 0x5C, 0xB4)
NAVY       = RGBColor(0x1E, 0x29, 0x3B)
SLATE      = RGBColor(0x33, 0x41, 0x55)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF0, 0xF5, 0xFA)
BORDER     = RGBColor(0xD1, 0xD5, 0xDB)
RED        = RGBColor(0xDC, 0x26, 0x26)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
ORANGE     = RGBColor(0xEA, 0x58, 0x0C)
GRAY       = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─── Helpers ──────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]  # blank layout


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, w, h, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_box(slide, left, top, w, h, text, fill_color=BLUE, font_color=WHITE,
            font_size=14, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE):
    shape = add_rect(slide, left, top, w, h, fill_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = align
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    shape.text_frame.margin_left = Pt(6)
    shape.text_frame.margin_right = Pt(6)
    shape.text_frame.margin_top = Pt(4)
    shape.text_frame.margin_bottom = Pt(4)
    return shape


def add_arrow_right(slide, left, top, w=Inches(0.5), h=Inches(0.3), color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_arrow_down(slide, left, top, w=Inches(0.3), h=Inches(0.4), color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color, text="", font_color=WHITE, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_text_box(slide, left, top, w, h, text, font_size=14, color=SLATE,
                 bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_bullets(slide, left, top, w, h, items, font_size=14, color=SLATE,
                bullet_char="\u2022", spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.space_before = Pt(2)
    return txBox


def add_sub_bullets(slide, left, top, w, h, items, font_size=12, color=GRAY,
                    bullet_char="\u25B8", spacing=Pt(3)):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"    {bullet_char}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
    return txBox


def slide_header(slide, title, subtitle=None, dark=False):
    """Add top blue bar + title + optional subtitle."""
    bar_color = NAVY if dark else BLUE
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), bar_color)
    add_text_box(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.6),
                 title, font_size=28, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.7), Inches(12), Inches(0.35),
                     subtitle, font_size=14, color=LIGHT_BLUE)
    # thin accent line
    add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), BLUE_DARK)


def slide_footer(slide, text="ElenzaIndia.com  |  Production Management System"):
    add_text_box(slide, Inches(0.5), Inches(7.05), Inches(12), Inches(0.35),
                 text, font_size=9, color=GRAY, align=PP_ALIGN.CENTER)


def section_title_slide(slide, section_num, section_title, icon_text=""):
    add_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(2.8), BLUE)
    if icon_text:
        add_circle(slide, Inches(5.8), Inches(1.2), Inches(1.6), BLUE_DARK,
                    icon_text, WHITE, 36)
    add_text_box(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
                 f"SECTION {section_num}", font_size=18, color=LIGHT_BLUE,
                 bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(1.2),
                 section_title, font_size=36, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.6), Inches(11), Inches(0.5),
                 "ElenzaIndia.com Production Management System",
                 font_size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


def two_col_slide(slide, title, left_title, left_items, right_title, right_items,
                  subtitle=None):
    slide_header(slide, title, subtitle)
    # Left card
    add_rect(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.2), LIGHT_BG,
             BORDER, 0.03)
    add_text_box(slide, Inches(0.8), Inches(1.55), Inches(5.2), Inches(0.45),
                 left_title, font_size=18, color=BLUE, bold=True)
    add_bullets(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(4.2),
                left_items, font_size=13)
    # Right card
    add_rect(slide, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5.2), LIGHT_BG,
             BORDER, 0.03)
    add_text_box(slide, Inches(7.1), Inches(1.55), Inches(5.2), Inches(0.45),
                 right_title, font_size=18, color=BLUE, bold=True)
    add_bullets(slide, Inches(7.1), Inches(2.1), Inches(5.2), Inches(4.2),
                right_items, font_size=13)
    slide_footer(slide)


def table_slide(slide, title, headers, rows, col_widths=None, subtitle=None):
    slide_header(slide, title, subtitle)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_w = sum(col_widths) if col_widths else Inches(12)
    left = (SLIDE_W - tbl_w) // 2
    top = Inches(1.5)
    row_h = Inches(0.45)
    tbl_h = row_h * n_rows
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, tbl_w, tbl_h)
    table = table_shape.table
    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            bg = LIGHT_BG if i % 2 == 0 else WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = SLATE
                p.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Set column widths
    if col_widths:
        for j, w in enumerate(col_widths):
            table.columns[j].width = w
    slide_footer(slide)


def flow_slide(slide, title, steps, subtitle=None):
    """Horizontal flow with boxes and arrows."""
    slide_header(slide, title, subtitle)
    n = len(steps)
    box_w = Inches(1.6)
    gap = Inches(0.25)
    arrow_w = Inches(0.4)
    total = n * box_w + (n - 1) * (gap + arrow_w)
    start_x = (SLIDE_W - total) // 2
    y_box = Inches(2.8)
    y_arrow = y_box + box_w // 2 - Inches(0.15)
    colors = [BLUE, BLUE_DARK, NAVY, BLUE, BLUE_DARK, NAVY, BLUE, BLUE_DARK, NAVY, BLUE]
    for i, step in enumerate(steps):
        x = start_x + i * (box_w + gap + arrow_w)
        c = colors[i % len(colors)]
        add_box(slide, x, y_box, box_w, box_w, step, c, WHITE, 12, True)
        if i < n - 1:
            add_arrow_right(slide, x + box_w + gap // 2, y_arrow, arrow_w,
                            Inches(0.25), BLUE_DARK)
    slide_footer(slide)


def icon_card_row(slide, cards, y=Inches(1.5)):
    """Row of icon cards: list of (icon_text, label, desc)."""
    n = len(cards)
    card_w = Inches(2.2)
    gap = Inches(0.3)
    total = n * card_w + (n - 1) * gap
    start_x = (SLIDE_W - total) // 2
    for i, (icon, label, desc) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rect(slide, x, y, card_w, Inches(4.5), LIGHT_BG, BORDER, 0.05)
        add_circle(slide, x + card_w // 2 - Inches(0.45), y + Inches(0.3),
                   Inches(0.9), BLUE, icon, WHITE, 28)
        add_text_box(slide, x + Inches(0.15), y + Inches(1.4), card_w - Inches(0.3),
                     Inches(0.4), label, font_size=13, color=NAVY, bold=True,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.15), y + Inches(1.85), card_w - Inches(0.3),
                     Inches(2.4), desc, font_size=10, color=SLATE,
                     align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE BUILDERS — organised by section
# ═══════════════════════════════════════════════════════════════════

def build_cover(prs):
    # Slide 1 — Cover
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(7.5), NAVY)
    add_rect(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.5), BLUE, None, 0.02)
    add_text_box(s, Inches(1.5), Inches(1.9), Inches(10.3), Inches(0.6),
                 "ELENZAINIDA.COM", font_size=18, color=LIGHT_BLUE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.5),
                 "Production Management\nSystem", font_size=44, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5),
                 "Modular Interiors B2B Operations Platform",
                 font_size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(5.5), Inches(5.0), Inches(2.3), Inches(0.04), WHITE)
    add_text_box(s, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.4),
                 "Operations Overview  |  Pitch Deck",
                 font_size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    dt = datetime.date.today().strftime("%B %Y")
    add_text_box(s, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.4),
                 dt, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # Slide 2 — Agenda
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Agenda", "What We Will Cover Today")
    items = [
        ("01", "The Problem — Why Factories Struggle"),
        ("02", "Our Solution — Elenza PMS Overview"),
        ("03", "Core Modules Deep Dive (13 Modules)"),
        ("04", "Production Flow & Machine Tracking"),
        ("05", "Planner, Reports & Audit"),
        ("06", "Technology & Architecture"),
        ("07", "Value Proposition & ROI"),
        ("08", "Implementation & Training"),
        ("09", "Future Roadmap"),
        ("10", "Q&A"),
    ]
    y = Inches(1.5)
    for num, label in items:
        add_circle(s, Inches(1.2), y, Inches(0.55), BLUE, num, WHITE, 16)
        add_text_box(s, Inches(2.0), y + Inches(0.05), Inches(10), Inches(0.45),
                     label, font_size=16, color=NAVY, bold=True)
        y += Inches(0.55)
    slide_footer(s)

    # Slide 3 — Why Elenza PMS
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Why Elenza PMS?", "Built for Modular Interior Manufacturing")
    cards = [
        ("\u2699", "Purpose-Built", "Designed specifically\nfor modular interior\nB2B operations"),
        ("\u260F", "Data-Entry First", "Fast forms, searchable\ndropdowns, minimal\nmanual typing"),
        ("\u2611", "Role-Based", "Every user sees only\nwhat they need —\nnothing more"),
        ("\u2618", "Audit Trail", "Every action tracked.\nFull order lifecycle\nvisibility"),
        ("\u2696", "Configurable", "Admin manages machines,\nsequences, masters —\nno developer needed"),
        ("\u26A1", "Zero Dashboard Clutter", "Factory-floor focused.\nTables, not charts."),
    ]
    icon_card_row(s, cards, y=Inches(1.5))
    slide_footer(s)


def build_problem(prs):
    # Slide 4 — Section title
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "01", "The Problem", "\u26A0")

    # Slide 5 — Pain Points Overview
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Current Manufacturing Challenges",
                 "What Factory Owners Face Every Day")
    pains = [
        ("\u2717", "Lost Orders", "No single system to track\nwhere each order is right now"),
        ("\u2717", "Manual Chaos", "Paper registers, WhatsApp\nmessages, phone calls"),
        ("\u2717", "No Audit Trail", "Cannot trace who changed\nwhat and when"),
        ("\u2717", "Permission Leaks", "Everyone sees everything\nor nothing — no middle ground"),
        ("\u2717", "No Rejection Tracking", "Rejected parts vanish\nwithout traceability"),
        ("\u2717", "Dispatch Gaps", "Packing done but dispatch\nstatus unknown"),
    ]
    icon_card_row(s, pains, y=Inches(1.5))
    slide_footer(s)

    # Slide 6 — Order Lifecycle Without PMS
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Order Lifecycle — Without a PMS",
                 "Fragmented, Untraceable, Error-Prone")
    steps = ["Enquiry\nReceived", "Quotation\nSent", "Order\nConfirmed", "Materials\nOrdered",
             "Production\nStart", "??? Where\nis it?", "Dispatch\n?"]
    flow_slide(s, "", steps)
    add_text_box(s, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
                 "At every stage: uncertainty, phone calls, manual checks, lost information",
                 font_size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    slide_footer(s)

    # Slide 7 — Impact of No System
    s = prs.slides.add_slide(BLANK)
    two_col_slide(s,
        "Impact of Not Having a System",
        "Operational Impact",
        [
            "Orders get lost between departments",
            "Production bottlenecks go unnoticed",
            "Rejections require phone follow-ups",
            "Dispatch delays cause customer complaints",
            "No visibility into machine workload",
            "Partial completions are hard to track",
        ],
        "Business Impact",
        [
            "Missed delivery deadlines",
            "Customer trust erosion",
            "Revenue leakage from lost orders",
            "Exvertime due to rework without data",
            "Management flying blind on status",
            "No data for business decisions",
        ])

    # Slide 8 — What You Need
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "What a Modular Interior Factory Actually Needs",
                 "Simple, Fast, Reliable")
    needs = [
        ("1", "Fast Data Entry", "Quick quotation & order capture\nwith searchable dropdowns"),
        ("2", "Order Tracking", "Know exactly where every\norder is at any moment"),
        ("3", "Machine Visibility", "Each machine station has its\nown queue and actions"),
        ("4", "Production Intelligence", "Planner sees all active orders,\npriorities, and stages"),
        ("5", "Reporting", "Search, filter, export to Excel\n— all without dashboards"),
        ("6", "Audit Trail", "Every change logged with\nuser, timestamp, and remarks"),
    ]
    icon_card_row(s, needs, y= Inches(1.5))
    slide_footer(s)

    # Slide 9 — Existing Solutions Fail
    s = prs.slides.add_slide(BLANK)
    two_col_slide(s,
        "Why Existing Solutions Fall Short",
        "Generic ERP Systems",
        [
            "Too complex for factory floor",
            "Expensive licensing & maintenance",
            "Requires dedicated IT team",
            "Dashboard-heavy, data-entry-light",
            "Not designed for production flow",
            "Overkill for modular interiors",
        ],
        "Custom SaaS Products",
        [
            "Monthly recurring costs add up",
            "Data locked in vendor servers",
            "Internet dependency in factories",
            "Features you don't use, missing ones you need",
            "Vendor support may be slow",
            "Cannot customize for local workflows",
        ])


def build_solution(prs):
    # Slide 10 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "02", "The Solution", "\u2714")

    # Slide 11 — Elenza PMS Intro
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    add_rect(s, Inches(1), Inches(1.5), Inches(11.3), Inches(4.8), BLUE, None, 0.02)
    add_text_box(s, Inches(1.5), Inches(1.8), Inches(10.3), Inches(0.5),
                 "INTRODUCING", font_size=16, color=LIGHT_BLUE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.0),
                 "Elenza PMS", font_size=48, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.5), Inches(3.4), Inches(10.3), Inches(0.6),
                 "Production Management System", font_size=22, color=LIGHT_BLUE,
                 align=PP_ALIGN.CENTER)
    add_rect(s, Inches(5.2), Inches(4.2), Inches(2.9), Inches(0.04), WHITE)
    add_text_box(s, Inches(1.5), Inches(4.5), Inches(10.3), Inches(1.0),
                 "A minimal, data-entry-first system built\n"
                 "specifically for modular interior B2B operations",
                 font_size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.4),
                 "One System  |  Every Order  |  Every Machine  |  Every Stage",
                 font_size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Slide 12 — Vision
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Our Vision", "One Simple System for Factory & Office Teams")
    goals = [
        ("Minimize\nManual Entry", "Searchable dropdowns\nreduce typing by 60%+"),
        ("Track Every\nOrder", "One current stage\nvisible at all times"),
        ("Enforce\nPermissions", "Role-based access\nevery endpoint"),
        ("Predictable\nFlow", "Work moves cleanly\nbetween machines"),
        ("Full Audit\nTrail", "Every action logged\nwith user & time"),
        ("Exportable\nReports", "Search, filter, sort\nExcel-ready data"),
    ]
    y = Inches(1.5)
    for i, (title, desc) in enumerate(goals):
        col = i % 3
        row = i // 3
        x = Inches(0.7) + col * Inches(4.2)
        yy = y + row * Inches(2.6)
        add_rect(s, x, yy, Inches(3.8), Inches(2.2), LIGHT_BG, BORDER, 0.05)
        add_text_box(s, x + Inches(0.2), yy + Inches(0.25), Inches(3.4), Inches(0.7),
                     title, font_size=16, color=BLUE, bold=True,
                     align=PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.2), yy + Inches(1.1), Inches(3.4), Inches(0.9),
                     desc, font_size=12, color=SLATE, align=PP_ALIGN.CENTER)
    slide_footer(s)

    # Slide 13 — Key Differentiators
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Key Differentiators", "What Sets Elenza PMS Apart")
    diff = [
        "\u2714  Factory-floor-first UI — not a dashboard analytics tool",
        "\u2714  Machine-wise production tracking with station-level login",
        "\u2714  Partial completion visible at both current and next station",
        "\u2714  Rejection moves work backward with mandatory reason",
        "\u2714  Packing and Dispatch as dedicated operational portals",
        "\u2714  Planner workspace with priority, machine grouping, and Excel export",
        "\u2714  Configurable production sequence — admin manages without developer",
        "\u2714  Every master (customer type, order type, machines) admin-editable",
        "\u2714  Role + station authorization enforced server-side on every action",
        "\u2714  Lightweight stack — HTML/CSS/JS, C# ASP.NET, Microsoft Access",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                diff, font_size=15, color=SLATE)
    slide_footer(s)


def build_modules_overview(prs):
    # Slide 14 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "03", "Core Modules Deep Dive", "\u2699")

    # Slide 15 — Module Map
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "13 Core Modules", "Complete Operational Coverage")
    modules = [
        "\u25A3 Login & User Roles",
        "\u25A3 Dealer Management",
        "\u25A3 Quotation Entry",
        "\u25A3 Order Confirmation",
        "\u25A3 Optimisation",
        "\u25A3 Procurement",
        "\u25A3 Production Tracking",
        "\u25A3 Packing",
        "\u25A3 Dispatch",
        "\u25A3 Production Planner",
        "\u25A3 Reports & Analytics",
        "\u25A3 Master Settings",
        "\u25A3 Audit Trail",
    ]
    cols = 3
    for i, mod in enumerate(modules):
        col = i % cols
        row = i // cols
        x = Inches(0.6) + col * Inches(4.2)
        y = Inches(1.5) + row * Inches(1.2)
        bg = BLUE if i % 2 == 0 else BLUE_DARK
        add_box(s, x, y, Inches(3.8), Inches(0.95), mod, bg, WHITE, 14, True)
    slide_footer(s)

    # Slide 16 — Module Flow
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module Interconnection", "How Data Flows Between Modules")
    flow_steps = [
        "Dealer\nEntry",
        "Quotation\nEntry",
        "Order\nConfirm",
        "Optimise",
        "Procure",
        "Produce",
        "Pack",
        "Dispatch",
    ]
    flow_slide(s, "", flow_steps)
    add_text_box(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.8),
                 "Each module feeds the next. Orders move forward on completion.\n"
                 "Rejection at any stage sends work backward with mandatory reason.",
                 font_size=13, color=SLATE, align=PP_ALIGN.CENTER)
    slide_footer(s)


def build_login_module(prs):
    # Slide 17
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module 01 — Login & User Roles",
                 "Every User Has a Purpose-Built Landing")
    steps = [
        "Enter\nUsername",
        "Enter\nPassword",
        "System\nIdentifies\nRole",
        "Routes to\nCorrect\nSection",
        "Session\nTracks\nActivity",
    ]
    flow_slide(s, "", steps)
    slide_footer(s)

    # Slide 18 — User Roles Table
    s = prs.slides.add_slide(BLANK)
    table_slide(s, "User Roles & Permissions",
        ["Role", "Access Level", "Key Capabilities"],
        [
            ["Admin / Super User", "Full", "Manage all masters, users, sequences, data, reports, audit"],
            ["Data Entry", "Entry", "Add dealers, quotations, confirmations, optimisation, procurement"],
            ["Quotation User", "Entry", "Create and manage quotations and dealer interactions"],
            ["Marketing User", "Scoped", "View quotation and dealer activity for their scope"],
            ["Optimisation User", "Entry", "Record board optimisation and raw material details"],
            ["Procurement User", "Entry", "Manage POs, vendor selection, material receipt tracking"],
            ["Production Planner", "Planner", "View all active orders, assign priority, machine grouping"],
            ["Machine User", "Station", "See only assigned station orders, update status"],
            ["Packing User", "Station", "Record packed boxes, balance, complete to dispatch"],
            ["Dispatch User", "Station", "Manage dispatch readiness, transport, completion"],
            ["Management", "Read-Only", "View reports, order status, lifecycle — no edits"],
        ],
        [Inches(2.5), Inches(2), Inches(7.5)])

    # Slide 19 — Role-Based Access Visual
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Role-Based Access Control", "Server-Side Enforcement on Every Endpoint")
    roles = [
        ("Admin", "Everything", BLUE),
        ("Data Entry", "Dealers, Quotes,\nOrders", BLUE_DARK),
        ("Machine", "Assigned\nStation Only", NAVY),
        ("Packing", "Packing\nPortal", BLUE),
        ("Dispatch", "Dispatch\nQueue", BLUE_DARK),
        ("Management", "Reports\nOnly", NAVY),
        ("Planner", "Planning\nWorkspace", BLUE),
    ]
    x = Inches(0.5)
    for role, access, color in roles:
        add_box(s, x, Inches(2.0), Inches(1.6), Inches(1.2), role, color, WHITE, 13, True)
        add_box(s, x, Inches(3.4), Inches(1.6), Inches(1.0), access, LIGHT_BG, SLATE, 11, False)
        x += Inches(1.82)
    add_text_box(s, Inches(0.5), Inches(5.0), Inches(12), Inches(0.6),
                 "\u26A0  Machine users see ONLY their assigned station orders\n"
                 "     Admin has full override for corrections and management",
                 font_size=13, color=NAVY, bold=True)
    slide_footer(s)


def build_dealer_module(prs):
    # Slide 20
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module 02 — Dealer Data Entry",
                 "Complete Dealer Profile Management")
    fields = [
        "\u2022  Dealer ID (auto-generated)",
        "\u2022  Dealer Name",
        "\u2022  Contact Person",
        "\u2022  Mobile Number",
        "\u2022  WhatsApp Number",
        "\u2022  Email Address",
        "\u2022  Company Name",
        "\u2022  City & Area",
        "\u2022  GST Number",
        "\u2022  Full Address",
        "\u2022  Dealer Type",
        "\u2022  Active / Inactive Status",
        "\u2022  Remarks",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.5),
                fields, font_size=14, color=SLATE)
    rules = [
        "Dealer name is searchable in quotation entry",
        "Dropdown auto-suggests as you type",
        "Duplicate mobile or GST triggers warning",
        "Admin can edit and manage all dealer details",
        "Inactive dealers hidden from new entries",
        "Dealer customer type can be changed by Admin",
    ]
    add_bullets(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.5),
                rules, font_size=13, color=NAVY, bullet_char="\u2714")
    slide_footer(s)

    # Slide 21 — Dealer Customer Type Change
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Dealer Customer Type Change",
                 "Admin Tool — No Database Script Required")
    steps = [
        "Admin\nSelects\nDealer",
        "Chooses\nNew Customer\nType",
        "System\nValidates\nAgainst Master",
        "Updates\nDealer &\nLinked Orders",
        "Audit\nLog\nRecorded",
    ]
    flow_slide(s, "", steps)
    add_text_box(s, Inches(0.5), Inches(5.5), Inches(12), Inches(1.0),
                 "\u2714  Admin selects dealer from searchable dropdown\n"
                 "\u2714  Customer type values come from Customer Type Master\n"
                 "\u2714  Both dealer record and all linked orders are updated\n"
                 "\u2714  Full audit trail captured",
                 font_size=12, color=SLATE)
    slide_footer(s)


def build_quotation_module(prs):
    # Slide 22
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module 03 — Quotation Data Entry",
                 "Capture Enquiry & Quotation Details")
    fields = [
        "\u2022  Quotation Number (auto-generated)",
        "\u2022  Quotation Date",
        "\u2022  Dealer Name (dropdown)",
        "\u2022  Customer Name",
        "\u2022  Customer Type (from master)",
        "\u2022  Order Type (from master)",
        "\u2022  Main Order & Sub Order",
        "\u2022  Order Number (unique)",
        "\u2022  Site / Project Name",
        "\u2022  Location",
        "\u2022  Approximate Value",
        "\u2022  Expected Confirmation Date",
        "\u2022  Remarks",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.5),
                fields, font_size=14, color=SLATE)
    rules = [
        "Order number must be unique",
        "All dropdowns support type-and-search",
        "Customer type comes from Customer Type Master",
        "Main order and sub order are identifiable",
        "Created by and timestamp auto-captured",
        "Quotation register available for review",
        "Delete action available with confirmation",
    ]
    add_bullets(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.5),
                rules, font_size=13, color=NAVY, bullet_char="\u2714")
    slide_footer(s)

    # Slide 23 — Quotation Register
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Quotation Register",
                 "Searchable, Filterable Order Listing")
    features = [
        ("\u2315", "Search", "Real-time search across\nall quotation fields"),
        ("\u2261", "Filter", "By date range, dealer,\ncustomer type, order type"),
        ("\u21C5", "Sort", "Click any column header\nto sort ascending/descending"),
        ("\u21E9", "Export", "Download filtered results\nto Excel-compatible CSV"),
        ("\u2716", "Delete", "Delete / More action with\nconfirmation protection"),
    ]
    icon_card_row(s, features, y=Inches(1.5))
    slide_footer(s)


def build_confirmation_module(prs):
    # Slide 24
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module 04 — Order Confirmation",
                 "Mark Quotation as Confirmed")
    steps = [
        "Select\nUnconfirmed\nOrder",
        "Set\nConfirmation\nDate",
        "Record\nConfirmed\nBy",
        "Add\nRemarks",
        "Status\nChanges to\nConfirmed",
    ]
    flow_slide(s, "", steps)
    add_text_box(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.8),
                 "\u2714  Dropdown shows only not-yet-confirmed orders\n"
                 "\u2714  Once confirmed, order becomes eligible for Optimisation\n"
                 "\u2714  Confirmation date and user are auto-recorded",
                 font_size=13, color=SLATE, align=PP_ALIGN.CENTER)
    slide_footer(s)

    # Slide 25 — Confirmation Rules
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Order Confirmation — Rules & Behaviour")
    rules = [
        "Order dropdown filters: status = Created (not yet confirmed)",
        "Confirmation date is mandatory",
        "Confirmed by field auto-populates from session user",
        "Remarks field available for any notes",
        "Once confirmed: order enters production planning pipeline",
        "Order is no longer available in confirmation dropdown",
        "Optimisation module can now pick this order",
        "Audit log records the confirmation event with timestamp",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                rules, font_size=15, color=SLATE, bullet_char="\u2714")
    slide_footer(s)


def build_optimisation_module(prs):
    # Slide 26
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module 05 — Optimisation",
                 "Board Optimisation & Raw Material Planning")
    fields = [
        "\u2022  Order Number (dropdown — confirmed, not optimised)",
        "\u2022  Optimisation Date",
        "\u2022  Number of Boards (numeric)",
        "\u2022  RM (Raw Material) Details",
        "\u2022  Optimisation Done By",
        "\u2022  Remarks",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.5),
                fields, font_size=15, color=SLATE)
    rules = [
        "Order dropdown shows only confirmed but not optimised orders",
        "Board count must be numeric",
        "After optimisation: order enters procurement pipeline",
        "Order no longer appears in optimisation dropdown",
        "Audit trail captures optimisation event",
        "Board quantity feeds into planner visibility",
    ]
    add_bullets(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.5),
                rules, font_size=14, color=NAVY, bullet_char="\u2714")
    slide_footer(s)

    # Slide 27 — Optimisation Value
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Optimisation — Why It Matters",
                 "Connecting Planning to Procurement")
    values = [
        ("Board\nCount", "Exact number of boards\nneeded for production"),
        ("RM\nDetails", "Raw material specifications\nfor procurement team"),
        ("Traceability", "Who optimised, when,\nand for which order"),
        ("Pipeline\nTrigger", "Unlocks procurement\nmodule for the order"),
        ("Planner\nVisibility", "Board qty visible in\nplanner workspace"),
    ]
    x = Inches(0.4)
    for title, desc in values:
        add_box(s, x, Inches(2.0), Inches(2.3), Inches(1.1), title, BLUE, WHITE, 13, True)
        add_text_box(s, x, Inches(3.3), Inches(2.3), Inches(1.2),
                     desc, font_size=11, color=SLATE, align=PP_ALIGN.CENTER)
        if x > Inches(0.5):
            add_arrow_right(s, x - Inches(0.5), Inches(2.4), Inches(0.4),
                            Inches(0.25), BLUE_DARK)
        x += Inches(2.6)
    slide_footer(s)


def build_procurement_module(prs):
    # Slide 28
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module 06 — Procurement",
                 "Purchase Order & Material Receiving Tracker")
    fields = [
        "\u2022  Order Number (dropdown)",
        "\u2022  PO Raised Date",
        "\u2022  PO Number (unique)",
        "\u2022  Vendor Name (dropdown)",
        "\u2022  Item Details",
        "\u2022  MRN (Material Received Note) Date",
        "\u2022  Procurement Status",
        "\u2022  Remarks",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.5),
                fields, font_size=15, color=SLATE)
    statuses = [
        "PO Pending — awaiting PO creation",
        "PO Raised — sent to vendor",
        "Partial Material Received — some items arrived",
        "Material Received — all items received",
        "Cancelled — order or PO cancelled",
    ]
    add_bullets(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.5),
                statuses, font_size=13, color=NAVY, bullet_char="\u25B8")
    slide_footer(s)

    # Slide 29 — Procurement Flow
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Procurement Flow",
                 "From PO to Material Receipt")
    steps = [
        "PO\nPending",
        "PO\nRaised",
        "Partial\nReceived",
        "Material\nReceived",
        "Ready for\nProduction",
    ]
    flow_slide(s, "", steps)
    add_text_box(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.8),
                 "\u2714  Multiple items may be added against one order\n"
                 "\u2714  PO number uniqueness enforced\n"
                 "\u2714  Material Received status makes order eligible for production",
                 font_size=13, color=SLATE, align=PP_ALIGN.CENTER)
    slide_footer(s)


def build_production_tracking(prs):
    # Slide 30 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "03A", "Production Tracking", "\u2699")

    # Slide 31 — Default Sequence
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Default Production Sequence",
                 "7-Station Factory Floor Flow")
    stations = [
        "Hot\nPress", "Cutting", "Edge-\nbanding", "Drilling",
        "QC", "Packing", "Dispatch",
    ]
    flow_slide(s, "", stations)
    add_text_box(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.8),
                 "Admin can: Add machines  |  Rename  |  Disable  |  Change sequence  |  Insert between",
                 font_size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    slide_footer(s)

    # Slide 32 — Machine Screen
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Machine Station Screen",
                 "What the Machine User Sees")
    fields = [
        "\u2022  Order Number",
        "\u2022  Dealer Name",
        "\u2022  Customer Name",
        "\u2022  Order Type",
        "\u2022  Main Order & Sub Order",
        "\u2022  Current Station",
        "\u2022  Previous Station",
        "\u2022  Next Station",
        "\u2022  Status Buttons: Completed | Partial | Rejected",
        "\u2022  Remarks Field",
        "\u2022  Last Updated Date & Time",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(6.0), Inches(5.5),
                fields, font_size=14, color=SLATE)
    add_rect(s, Inches(7.2), Inches(1.5), Inches(5.5), Inches(5.0), LIGHT_BG, BORDER, 0.03)
    add_text_box(s, Inches(7.5), Inches(1.7), Inches(5.0), Inches(0.4),
                 "Key Rules", font_size=16, color=BLUE, bold=True)
    rules = [
        "Machine user sees ONLY assigned station orders",
        "Cannot update orders from other stations",
        "Completed: order moves to next station",
        "Partial: visible at current AND next station",
        "Rejected: order moves back to previous",
        "First station rejection: admin correction queue",
        "Remarks mandatory for Partial and Rejected",
    ]
    add_bullets(s, Inches(7.5), Inches(2.2), Inches(5.0), Inches(4.0),
                rules, font_size=12, color=NAVY)
    slide_footer(s)

    # Slide 33 — Machine Admin Management
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Machine / Station Management",
                 "Admin Configurable Without Developer")
    capabilities = [
        ("\u2795", "Add", "Add new machine or\nstation to the sequence"),
        ("\u270F", "Rename", "Rename existing\nmachine or station"),
        ("\u2716", "Disable", "Disable station —\nremoved from new entries"),
        ("\u21C5", "Reorder", "Change production\nsequence order"),
        ("\u279C", "Insert", "Insert machine between\nexisting stations"),
    ]
    icon_card_row(s, capabilities, y= Inches(1.5))
    slide_footer(s)

    # Slide 34 — Production Movement: Completed
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Production Movement — Completed",
                 "Order Moves Forward to Next Station")
    add_rect(s, Inches(1), Inches(1.8), Inches(11.3), Inches(4.5), LIGHT_BG, BORDER, 0.02)
    add_text_box(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.5),
                 "When Machine User Marks an Order as COMPLETED:",
                 font_size=18, color=BLUE, bold=True)
    steps = [
        "Cutting\nStation",
        "COMPLETED",
        "Edgebanding\nStation",
    ]
    y = Inches(2.8)
    x = Inches(2.5)
    add_box(s, x, y, Inches(2.5), Inches(1.0), steps[0], NAVY, WHITE, 14, True)
    add_arrow_right(s, x + Inches(2.7), y + Inches(0.35), Inches(1.2),
                    Inches(0.3), GREEN)
    add_box(s, x + Inches(4.1), y, Inches(2.5), Inches(1.0), steps[2], BLUE, WHITE, 14, True)
    results = [
        "\u2714  Order DISAPPEARS from current machine login",
        "\u2714  Order APPEARS at next machine in sequence",
        "\u2714  If current station is Dispatch: order becomes Fully Completed",
        "\u2714  Fully completed orders do NOT appear in any machine login",
    ]
    add_bullets(s, Inches(1.5), Inches(4.2), Inches(10), Inches(2.0),
                results, font_size=14, color=NAVY, bullet_char="")
    slide_footer(s)

    # Slide 35 — Production Movement: Partial Completed
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Production Movement — Partial Completed",
                 "Order Visible at Both Current and Next Station")
    add_rect(s, Inches(1), Inches(1.8), Inches(11.3), Inches(4.5), LIGHT_BG, BORDER, 0.02)
    add_text_box(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.5),
                 "When Machine User Marks an Order as PARTIAL COMPLETED:",
                 font_size=18, color=BLUE, bold=True)
    y = Inches(2.8)
    x = Inches(1.5)
    add_box(s, x, y, Inches(3.0), Inches(1.0), "Cutting\n(VISIBLE)", NAVY, WHITE, 14, True)
    add_arrow_right(s, x + Inches(3.2), y + Inches(0.35), Inches(1.0),
                    Inches(0.25), ORANGE)
    add_box(s, x + Inches(4.4), y, Inches(3.0), Inches(1.0),
            "Edgebanding\n(VISIBLE)", BLUE, WHITE, 14, True)
    results = [
        "\u2714  Order REMAINS visible in current machine",
        "\u2714  Order ALSO appears in next machine",
        "\u2714  Remarks are MANDATORY — must explain partial status",
        "\u2714  Planner shows originating stage as Current Stage until downstream movement",
    ]
    add_bullets(s, Inches(1.5), Inches(4.2), Inches(10), Inches(2.0),
                results, font_size=14, color=NAVY, bullet_char="")
    slide_footer(s)

    # Slide 36 — Production Movement: Rejected
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Production Movement — Rejected",
                 "Order Moves Backward to Previous Station")
    add_rect(s, Inches(1), Inches(1.8), Inches(11.3), Inches(4.5), LIGHT_BG, BORDER, 0.02)
    add_text_box(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.5),
                 "When Machine User Marks an Order as REJECTED:",
                 font_size=18, color=BLUE, bold=True)
    y = Inches(2.8)
    x = Inches(4.5)
    add_box(s, x, y, Inches(3.0), Inches(1.0), "Edgebanding\nStation", BLUE, WHITE, 14, True)
    add_arrow_right(s, x - Inches(1.5), y + Inches(0.35), Inches(1.0),
                    Inches(0.25), RED)
    add_box(s, x - Inches(3.5), y, Inches(3.0), Inches(1.0), "Cutting\nStation", NAVY, WHITE, 14, True)
    results = [
        "\u2714  Order DISAPPEARS from current machine",
        "\u2714  Order MOVES BACK to previous machine",
        "\u2714  Rejection reason is MANDATORY",
        "\u2714  First machine rejection: goes to Admin / Data Entry correction queue",
        "\u2714  Planner reapproval required before reintroduction into production",
    ]
    add_bullets(s, Inches(1.5), Inches(4.2), Inches(10), Inches(2.0),
                results, font_size=14, color=NAVY, bullet_char="")
    slide_footer(s)


def build_packing(prs):
    # Slide 37 — Packing Module
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module 07 — Packing Portal",
                 "Dedicated Packing Workspace")
    features = [
        ("\u2611", "Box Count", "Record packed boxes\nper order"),
        ("\u2610", "Balance", "Track remaining\nbalance quantity"),
        ("\u2714", "Complete", "Zero balance =\nmove to Dispatch"),
        ("\u231B", "History", "Persistent packing\nhistory after dispatch"),
        ("\u26D4", "No Duplicates", "Duplicate-submit\nprotection active"),
    ]
    icon_card_row(s, features, y= Inches(1.5))
    slide_footer(s)

    # Slide 38 — Packing Flow
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Packing Operation Flow",
                 "From Packing Queue to Dispatch Readiness")
    steps = [
        "Order\nArrives at\nPacking",
        "Enter\nPacked\nBoxes",
        "Enter\nBalance\nBoxes",
        "Save\nUpdate",
        "Balance\nZero?\nYES",
        "Move to\nDispatch",
    ]
    flow_slide(s, "", steps)
    add_text_box(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.8),
                 "\u2714  Packing station accepts both 'Packing' and 'Packed' labels\n"
                 "\u2714  Completed orders leave active queue but remain in history\n"
                 "\u2714  History shows packed boxes, balance, station, and action date",
                 font_size=12, color=SLATE, align=PP_ALIGN.CENTER)
    slide_footer(s)

    # Slide 39 — Packing Compatibility
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Packing Station Compatibility",
                 "Packing vs Packed — Handled Transparently")
    rows = [
        ["Station Name in Database", "Packed", "Legacy naming convention"],
        ["Portal Label", "Packing / Packed", "Both accepted"],
        ["User Assignment", "packing.user", "Assigned to 'Packed' station"],
        ["Authorization", "Role + Station", "Server-side check"],
        ["Save Behavior", "Submits assigned station", "No hardcoded value"],
        ["History", "Persistent", "Visible after dispatch"],
        ["Duplicate Protection", "Active", "Prevents double-click saves"],
    ]
    table_slide(s, "", ["Parameter", "Value", "Note"], rows,
                [Inches(3), Inches(3), Inches(6)])
    slide_footer(s)


def build_dispatch(prs):
    # Slide 40
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module 08 — Dispatch",
                 "Mandatory Stage After Packing")
    fields = [
        "\u2022  Order Number",
        "\u2022  Dealer Name",
        "\u2022  Customer Name",
        "\u2022  Order Type",
        "\u2022  Packing Completed Date",
        "\u2022  Dispatch Date",
        "\u2022  Vehicle / Transport Details",
        "\u2022  Dispatch Remarks",
        "\u2022  Dispatch Status",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.5),
                fields, font_size=14, color=SLATE)
    statuses = [
        "Pending Dispatch — packed, awaiting dispatch action",
        "Partially Dispatched — some items dispatched",
        "Dispatched — fully dispatched, removed from active view",
        "Hold — order held, remarks mandatory",
    ]
    add_bullets(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.5),
                statuses, font_size=13, color=NAVY, bullet_char="\u25B8")
    slide_footer(s)

    # Slide 41 — Dispatch Rules
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Dispatch — Rules & Visibility",
                 "Separate Login, Separate Queue")
    rules = [
        "Dispatch user has a separate login and dedicated queue",
        "Dispatch completed order disappears from Dispatch login",
        "Dispatched orders appear ONLY in reports",
        "Partial dispatch remains visible in Dispatch",
        "Dispatch remarks mandatory for Hold or Partially Dispatched",
        "Packed, Dispatch, and Dispatched rows NOT in Planner Priority Desk",
        "Dispatched orders NOT in production planner workspace",
        "Current Stage becomes 'Dispatched' after full dispatch",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                rules, font_size=14, color=SLATE, bullet_char="\u2714")
    slide_footer(s)


def build_planner(prs):
    # Slide 42 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "04", "Production Planner", "\u2611")

    # Slide 43 — Planner Workspace
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module 09 — Production Planner Workspace",
                 "Separate Page for Planner User")
    tabs = [
        ("\u25A3", "All Orders", "Complete list of all\nactive planning orders\nfrom confirmed to packing"),
        ("\u25A3", "Machine Wise", "Orders grouped by\nmachine/station with\nown export per table"),
        ("\u25A3", "Priority Desk", "Priority-focused view\nwith High/Med/Low\nsorting"),
    ]
    x = Inches(1.2)
    for icon, title, desc in tabs:
        add_rect(s, x, Inches(1.8), Inches(3.3), Inches(4.0), LIGHT_BG, BORDER, 0.04)
        add_circle(s, x + Inches(1.1), Inches(2.1), Inches(1.0), BLUE, icon, WHITE, 32)
        add_text_box(s, x + Inches(0.2), Inches(3.3), Inches(2.9), Inches(0.5),
                     title, font_size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.2), Inches(3.9), Inches(2.9), Inches(1.5),
                     desc, font_size=12, color=SLATE, align=PP_ALIGN.CENTER)
        x += Inches(3.7)
    slide_footer(s)

    # Slide 44 — Planner Visibility Rules
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Planner Visibility Rules",
                 "One Current Stage — Consistent Across All Views")
    rules = [
        "Show active planning orders from Confirmed stage up to Packing stage",
        "Do NOT show Dispatched orders — they are removed from planner",
        "Show only ONE 'Current Stage' per order at all times",
        "Do NOT show 'Visible Stations' as a planner display field",
        "Visible station routing is used internally, NOT as planner stage truth",
        "Partial completion: Current Stage remains originating stage",
        "Packed overrides earlier production stages until dispatch",
        "Dispatched: order exits planner workspace entirely",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                rules, font_size=14, color=SLATE, bullet_char="\u2714")
    slide_footer(s)

    # Slide 45 — Planner Fields
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Planner Workspace Fields",
                 "What the Planner Sees")
    fields = [
        ["Order Number", "Unique order identifier", "Text"],
        ["Dealer", "Dealer name from master", "Dropdown-linked"],
        ["Customer", "Customer name", "Text"],
        ["Order Type", "Type from master", "Dropdown-linked"],
        ["Current Stage", "One latest valid stage", "Auto-calculated"],
        ["Priority", "High / Medium / Low / Blank", "Editable from table"],
        ["EDD", "Expected Delivery Date", "Date"],
        ["Confirmation Date", "When order was confirmed", "Auto-captured"],
        ["Panel Qty", "Number of panels", "From optimisation"],
        ["Board Qty", "Number of boards", "From optimisation"],
    ]
    table_slide(s, "", ["Field", "Description", "Source"], fields,
                [Inches(3), Inches(5.5), Inches(3.5)])
    slide_footer(s)

    # Slide 46 — Priority System
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Planner Priority System",
                 "Direct Assignment from Planning Tables")
    priorities = [
        ("HIGH", RED, "Urgent orders float\nabove all others"),
        ("MEDIUM", ORANGE, "Standard priority\nfor normal orders"),
        ("LOW", GREEN, "Lower priority\nwhen capacity allows"),
        ("BLANK", GRAY, "No priority set\n— default state"),
    ]
    x = Inches(1.0)
    for label, color, desc in priorities:
        add_box(s, x, Inches(2.0), Inches(2.5), Inches(1.2), label, color, WHITE, 20, True)
        add_text_box(s, x, Inches(3.4), Inches(2.5), Inches(1.2),
                     desc, font_size=12, color=SLATE, align=PP_ALIGN.CENTER)
        x += Inches(3.0)
    add_text_box(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.0),
                 "\u2714 Priority assignable directly from planner tables — not just displayed\n"
                 "\u2714 High sorts above Medium, Low, and blank during planner sorting\n"
                 "\u2714 Priority updates captured in audit trail",
                 font_size=13, color=NAVY)
    slide_footer(s)

    # Slide 47 — Planner Export
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Planner Excel Export",
                 "Filter-Aware, Machine-Specific Exports")
    exports = [
        ("\u2611", "All Orders Tab", "Full planner export\nrespecting current filters"),
        ("\u2611", "Machine Wise Tab", "Consolidated export\nof all machine tables"),
        ("\u2611", "Single Machine", "Per-station export\nfor individual machines"),
        ("\u2611", "Priority Desk", "Priority-filtered export\nwith sorting preserved"),
        ("\u2611", "WIP Dated Files", "Auto-named with\ncurrent date stamp"),
    ]
    icon_card_row(s, exports, y= Inches(1.5))
    slide_footer(s)

    # Slide 48 — Machine Wise Sorting
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Machine Wise Export — Sorting Logic",
                 "Confirmation Date Descending")
    add_rect(s, Inches(1), Inches(1.8), Inches(11.3), Inches(4.5), LIGHT_BG, BORDER, 0.02)
    add_text_box(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.5),
                 "Export Behaviour:", font_size=18, color=BLUE, bold=True)
    rules = [
        "\u2714  Each machine table supports its own independent Excel export",
        "\u2714  Machine-wise export sorted by Confirmation Date — descending (newest first)",
        "\u2714  Filters applied in the UI are respected in the export",
        "\u2714  Export filenames: 'Machine Wise WIP Orders Dated DD-MM-YYYY.csv'",
        "\u2714  Single machine export: '<Station> WIP Orders Dated DD-MM-YYYY.csv'",
        "\u2714  Consolidated export: 'WIP Orders Dated DD-MM-YYYY.xlsx'",
    ]
    add_bullets(s, Inches(1.5), Inches(2.8), Inches(10), Inches(3.5),
                rules, font_size=14, color=SLATE, bullet_char="")
    slide_footer(s)


def build_reports(prs):
    # Slide 49 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "05", "Reports & Analytics", "\u2611")

    # Slide 50 — Reports Overview
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module 10 — Reports Tab",
                 "Separate from Data Entry — Comprehensive Reporting")
    reports = [
        "1.  Dealer Report",
        "2.  Quotation Report",
        "3.  Confirmed Order Report",
        "4.  Optimisation Report",
        "5.  Procurement Report",
        "6.  Production Status Report",
        "7.  Machine-Wise Pending Report",
        "8.  Rejected Order Report",
        "9.  Partial Completed Order Report",
        "10. Completed Order Report",
        "11. Dispatch Report",
        "12. Order Lifecycle Report",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.5),
                reports, font_size=14, color=SLATE)
    features = [
        "Every report supports: Search, Sort, Filter",
        "Date range filter available",
        "Export to Excel-compatible format",
        "Status filter for production states",
        "Dealer filter for dealer-specific data",
        "Order type filter for category views",
        "Machine / station filter for floor data",
        "Pagination for large result sets",
    ]
    add_bullets(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.5),
                features, font_size=13, color=NAVY, bullet_char="\u2714")
    slide_footer(s)

    # Slide 51 — Report Filter Capabilities
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Report Filter Capabilities",
                 "Every Report Has Full Search & Filter")
    filters = [
        ("\u2315", "Search", "Free-text search across\nreport fields"),
        ("\u2261", "Date Range", "Filter by creation,\nconfirmation, or action date"),
        ("\u2611", "Status", "Pending, In Progress,\nCompleted, Rejected, etc."),
        ("\u260F", "Dealer", "Filter by specific\ndealer or dealer group"),
        ("\u25A3", "Order Type", "Laminate, Kitchen,\nWardrobe, etc."),
        ("\u2699", "Machine", "Filter by production\nstation or stage"),
    ]
    icon_card_row(s, filters, y= Inches(1.5))
    slide_footer(s)

    # Slide 52 — Order Lifecycle Report
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Order Lifecycle Report",
                 "Complete Order Journey — Start to Dispatch")
    flow_steps = [
        "Quotation\nCreated",
        "Order\nConfirmed",
        "Optimised",
        "Procurement\nDone",
        "Production\nStarted",
        "Packing\nComplete",
        "Dispatched",
    ]
    flow_slide(s, "", flow_steps)
    add_text_box(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.8),
                 "Lifecycle shows: every status change, timestamp, user, machine, and remarks\n"
                 "Packed and Dispatched stage changes visible in lifecycle history",
                 font_size=12, color=SLATE, align=PP_ALIGN.CENTER)
    slide_footer(s)

    # Slide 53 — Production Status Report
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Production Status Report",
                 "Real-Time Factory Floor Visibility")
    statuses = [
        ("Pending", GRAY, "Awaiting action\nat current station"),
        ("In Progress", BLUE, "Currently being\nworked on"),
        ("Completed", GREEN, "Done at this\nstation"),
        ("Partial", ORANGE, "Partially done,\nvisible at 2 stations"),
        ("Rejected", RED, "Sent backward\nfor correction"),
        ("Dispatched", NAVY, "Fully shipped\nand closed"),
    ]
    x = Inches(0.5)
    for label, color, desc in statuses:
        add_box(s, x, Inches(2.0), Inches(1.9), Inches(1.0), label, color, WHITE, 13, True)
        add_text_box(s, x, Inches(3.2), Inches(1.9), Inches(1.0),
                     desc, font_size=10, color=SLATE, align=PP_ALIGN.CENTER)
        x += Inches(2.1)
    slide_footer(s)


def build_masters(prs):
    # Slide 54 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "06", "Master Settings", "\u2699")

    # Slide 55 — Masters Overview
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module 11 — Master Settings",
                 "Admin-Managed Configuration — No Developer Needed")
    masters = [
        ("\u25A3", "Dealer Master", "Dealer details,\ntypes, and status"),
        ("\u25A3", "Customer Type", "EL, ADM, MCB,\nBRGWF, and custom"),
        ("\u25A3", "Order Type", "15 production\ncategories"),
        ("\u25A3", "Machine / Station", "7 default stations\n+ custom"),
        ("\u25A3", "User Master", "Users, roles,\nand assignments"),
        ("\u25A3", "Vendor Master", "Vendor details\nand categories"),
    ]
    icon_card_row(s, masters, y= Inches(1.5))
    slide_footer(s)

    # Slide 56 — Customer Type Master
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Customer Type Master",
                 "Fully Admin-Configurable")
    rows = [
        ["EL", "Default Type", "Active"],
        ["ADM", "Default Type", "Active"],
        ["MCB", "Default Type", "Active"],
        ["BRGWF", "Default Type", "Active"],
        ["KADIWA", "Custom Type", "Active"],
        ["Custom 1", "Add New", "Admin adds as needed"],
    ]
    table_slide(s, "", ["Customer Type", "Description", "Status"], rows,
                [Inches(3), Inches(4), Inches(5)])
    add_text_box(s, Inches(0.5), Inches(5.8), Inches(12), Inches(1.0),
                 "\u2714 Admin can Add, Edit, Disable, Reorder, or Remove values\n"
                 "\u2714 Disabled types hidden from new entries but preserved in history\n"
                 "\u2714 Customer type dropdown in Quotation comes from this master",
                 font_size=12, color=NAVY)
    slide_footer(s)

    # Slide 57 — Order Type Master
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Order Type Master",
                 "15 Modular Interior Production Categories")
    types = [
        "Laminate", "Carcase", "Membrane", "Veneer", "Alu Glass",
        "Acrylic", "PU", "Profile Shutter", "Glass Shutter", "Wardrobe",
        "Kitchen", "TV Unit", "Full Home", "Loose Furniture", "Other",
    ]
    cols = 5
    for i, t in enumerate(types):
        col = i % cols
        row = i // cols
        x = Inches(0.6) + col * Inches(2.5)
        y = Inches(1.5) + row * Inches(1.4)
        add_box(s, x, y, Inches(2.2), Inches(1.1), t, BLUE if i % 2 == 0 else BLUE_DARK,
                WHITE, 13, True)
    add_text_box(s, Inches(0.5), Inches(5.8), Inches(12), Inches(1.0),
                 "\u2714 Admin can Add, Edit, Disable, Reorder order types\n"
                 "\u2714 Disabled types hidden from new quotation entries\n"
                 "\u2714 Future additions require no developer assistance",
                 font_size=12, color=NAVY)
    slide_footer(s)

    # Slide 58 — Machine Master
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Machine / Station Master",
                 "7 Default Stations — Fully Configurable")
    stations = [
        ("1", "Hot Press", "Board pressing\nand lamination"),
        ("2", "Cutting", "Board cutting\nto size"),
        ("3", "Edgebanding", "Edge banding\napplication"),
        ("4", "Drilling", "Hole drilling\nfor assembly"),
        ("5", "QC", "Quality check\nand inspection"),
        ("6", "Packing", "Box packing\nand counting"),
        ("7", "Dispatch", "Final dispatch\nto customer"),
    ]
    x = Inches(0.3)
    for num, name, desc in stations:
        add_circle(s, x + Inches(0.55), Inches(1.5), Inches(0.7), BLUE, num, WHITE, 18)
        add_text_box(s, x, Inches(2.4), Inches(1.8), Inches(0.4),
                     name, font_size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(s, x, Inches(2.85), Inches(1.8), Inches(0.8),
                     desc, font_size=10, color=SLATE, align=PP_ALIGN.CENTER)
        x += Inches(1.85)
    slide_footer(s)

    # Slide 59 — User Master
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "User Master",
                 "Role-Based User Management")
    fields = [
        "\u2022  User Name (display name)",
        "\u2022  Login ID (unique identifier)",
        "\u2022  Password (hashed)",
        "\u2022  Role (Admin, Data Entry, Machine, etc.)",
        "\u2022  Assigned Machine / Station (for Machine Users)",
        "\u2022  Active / Inactive Status",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(4.0),
                fields, font_size=14, color=SLATE)
    rules = [
        "Disabled users cannot login",
        "Machine users are linked to specific station",
        "Admin can create, edit, and deactivate users",
        "Role determines accessible modules",
        "Station assignment enforced server-side",
        "Password reset available through Admin tools",
    ]
    add_bullets(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(4.0),
                rules, font_size=13, color=NAVY, bullet_char="\u2714")
    slide_footer(s)

    # Slide 60 — Vendor Master
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Vendor Master",
                 "Procurement Vendor Management")
    fields = [
        "\u2022  Vendor Name",
        "\u2022  Contact Person",
        "\u2022  Material Category",
        "\u2022  Remarks",
        "\u2022  Active / Inactive Status",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(4.0),
                fields, font_size=15, color=SLATE)
    rules = [
        "Vendor dropdown in Procurement comes from this master",
        "Inactive vendors hidden from new entries",
        "Admin can add, edit, and manage vendors",
        "Material category links to procurement items",
    ]
    add_bullets(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(4.0),
                rules, font_size=14, color=NAVY, bullet_char="\u2714")
    slide_footer(s)


def build_audit(prs):
    # Slide 61 — Audit Trail
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Module 12 — Audit Trail",
                 "Every Important Action Logged")
    tracked = [
        "\u2022  Created by (user ID)",
        "\u2022  Created date and time",
        "\u2022  Updated by (user ID)",
        "\u2022  Updated date and time",
        "\u2022  Previous status",
        "\u2022  New status",
        "\u2022  Remarks / reason",
        "\u2022  Machine or station",
        "\u2022  User login used",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.5),
                tracked, font_size=14, color=SLATE)
    uses = [
        "Order lifecycle report shows full history",
        "Planner priority updates captured",
        "Packed and Dispatched transitions logged",
        "Rejection reasons preserved",
        "Machine user actions with timestamps",
        "Admin corrections fully traceable",
        "Visible station routing does NOT override stage history",
        "History grouped to one latest row per order",
    ]
    add_bullets(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.5),
                uses, font_size=13, color=NAVY, bullet_char="\u2714")
    slide_footer(s)

    # Slide 62 — Audit Trail Table
    s = prs.slides.add_slide(BLANK)
    table_slide(s, "Audit Trail — What Gets Captured",
        ["Action", "Captured Data", "Visible In"],
        [
            ["Dealer Created", "User, timestamp, dealer data", "Dealer Report"],
            ["Quotation Saved", "User, timestamp, all fields", "Quotation Report"],
            ["Order Confirmed", "User, date, remarks", "Lifecycle Report"],
            ["Optimisation Done", "User, boards, RM details", "Optimisation Report"],
            ["PO Raised", "User, PO#, vendor, items", "Procurement Report"],
            ["Machine Completed", "User, station, timestamp", "Production Report"],
            ["Partial Completed", "User, station, remarks", "Production Report"],
            ["Rejected", "User, station, reason", "Rejected Report"],
            ["Packed", "User, boxes, balance, date", "Packing History"],
            ["Dispatched", "User, date, transport, remarks", "Dispatch Report"],
            ["Priority Changed", "User, old/new priority", "Planner Audit"],
            ["Customer Type Changed", "User, old/new type", "Master Audit"],
        ],
        [Inches(2.8), Inches(4.5), Inches(4.7)])


def build_technology(prs):
    # Slide 63 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "07", "Technology & Architecture", "\u2699")

    # Slide 64 — Architecture Diagram
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "System Architecture",
                 "Simple, Reliable, Proven Stack")
    # Browser box
    add_box(s, Inches(1), Inches(2.0), Inches(3.5), Inches(1.5),
            "BROWSER\nHTML / CSS / JavaScript", NAVY, WHITE, 14, True)
    # Arrow
    add_arrow_right(s, Inches(4.8), Inches(2.5), Inches(0.8), Inches(0.3), BLUE)
    # API box
    add_box(s, Inches(5.8), Inches(2.0), Inches(3.5), Inches(1.5),
            "ASP.NET 4.8\nC# api.ashx", BLUE, WHITE, 14, True)
    # Arrow
    add_arrow_right(s, Inches(9.6), Inches(2.5), Inches(0.8), Inches(0.3), BLUE)
    # DB box
    add_box(s, Inches(10.6), Inches(2.0), Inches(2.2), Inches(1.5),
            "Access\n.accdb", NAVY, WHITE, 14, True)
    # Labels
    add_text_box(s, Inches(1), Inches(3.7), Inches(3.5), Inches(0.5),
                 "Static pages served by IIS\nSession-based authentication",
                 font_size=10, color=SLATE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(5.8), Inches(3.7), Inches(3.5), Inches(0.5),
                 "Action-based JSON API\nRole + Station authorization",
                 font_size=10, color=SLATE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(10.6), Inches(3.7), Inches(2.2), Inches(0.5),
                 "OleDb connection\nFile-based persistence",
                 font_size=10, color=SLATE, align=PP_ALIGN.CENTER)
    # Bottom row
    techs = [
        ("\u2699", "HTML/CSS/JS", "Plain frontend,\nno framework"),
        ("\u2699", "C# ASP.NET 4.8", "Server-side logic,\nauthorization"),
        ("\u2699", "Microsoft Access", "Lightweight DB,\nno server install"),
        ("\u2699", "FTP Deploy", "Simple deployment,\nno CI/CD needed"),
        ("\u2699", "Session Auth", "Cookie-based,\nserver-enforced"),
    ]
    x = Inches(0.4)
    for icon, label, desc in techs:
        add_rect(s, x, Inches(4.6), Inches(2.3), Inches(2.0), LIGHT_BG, BORDER, 0.04)
        add_circle(s, x + Inches(0.7), Inches(4.75), Inches(0.7), BLUE, icon, WHITE, 22)
        add_text_box(s, x + Inches(0.1), Inches(5.55), Inches(2.1), Inches(0.4),
                     label, font_size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.1), Inches(5.95), Inches(2.1), Inches(0.6),
                     desc, font_size=10, color=SLATE, align=PP_ALIGN.CENTER)
        x += Inches(2.55)
    slide_footer(s)

    # Slide 65 — Technology Advantages
    s = prs.slides.add_slide(BLANK)
    two_col_slide(s,
        "Why This Technology Stack?",
        "Business Advantages",
        [
            "No monthly SaaS fees — own the system",
            "No internet dependency for core operations",
            "Familiar Microsoft ecosystem for IT teams",
            "Simple deployment — FTP upload",
            "Database file is portable and backupable",
            "No vendor lock-in — full source ownership",
        ],
        "Technical Advantages",
        [
            "Proven ASP.NET Framework 4.8 stability",
            "Lightweight OleDb — no database server needed",
            "Plain HTML/CSS/JS — no build pipeline",
            "Session-based auth — server-enforced",
            "File-based database — easy backup/restore",
            "IIS-compatible shared hosting available",
        ])

    # Slide 66 — API Structure
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "API Structure",
                 "Action-Based JSON Endpoints")
    endpoints = [
        ["session", "Check current session state", "GET"],
        ["login-init", "Get login page data", "GET"],
        ["login", "Authenticate user", "POST"],
        ["logout", "End session", "POST"],
        ["app-state", "Full application state after login", "GET"],
        ["history-state", "Order history for current user", "GET"],
        ["save-dealer", "Create or update dealer", "POST"],
        ["save-quotation", "Create or update quotation", "POST"],
        ["confirm-order", "Mark order as confirmed", "POST"],
        ["save-optimisation", "Record board optimisation", "POST"],
        ["save-procurement", "Record PO and material receipt", "POST"],
        ["update-production", "Machine status action", "POST"],
        ["save-packing", "Record packing counts", "POST"],
        ["save-dispatch", "Record dispatch details", "POST"],
        ["save-priority", "Update planner priority", "POST"],
    ]
    table_slide(s, "", ["Endpoint", "Description", "Method"], endpoints,
                [Inches(3), Inches(6.5), Inches(2.5)])
    slide_footer(s)

    # Slide 67 — Security Features
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Security Features",
                 "Server-Side Authorization on Every Action")
    features = [
        ("\u2611", "Session Auth", "Cookie-based session\nwith user ID reload"),
        ("\u2611", "Role Checks", "Every endpoint verifies\nuser role"),
        ("\u2611", "Station Auth", "Machine users validated\nagainst assigned station"),
        ("\u2611", "Parameterized", "OleDb commands use\nparameterized queries"),
        ("\u2611", "Password Hash", "Passwords stored\nas hashes"),
        ("\u2611", "Audit Trail", "Every action logged\nwith user and timestamp"),
    ]
    icon_card_row(s, features, y= Inches(1.5))
    slide_footer(s)


def build_value_proposition(prs):
    # Slide 68 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "08", "Value Proposition & ROI", "\u2605")

    # Slide 69 — ROI
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Return on Investment",
                 "Measurable Efficiency Gains")
    metrics = [
        ("60%+", "Faster Data Entry", "Searchable dropdowns\neliminate manual typing"),
        ("100%", "Order Visibility", "Every order tracked\nfrom quote to dispatch"),
        ("0", "Lost Orders", "Complete lifecycle\nwith audit trail"),
        ("80%", "Fewer Phone Calls", "Planner sees all\nstatus in real time"),
        ("100%", "Machine Accountability", "Each station has\nits own queue"),
        ("100%", "Dispatch Accuracy", "Packing + dispatch\nportal tracking"),
    ]
    x = Inches(0.4)
    for value, label, desc in metrics:
        add_box(s, x, Inches(1.6), Inches(1.9), Inches(1.3), value, GREEN, WHITE, 28, True)
        add_text_box(s, x, Inches(3.1), Inches(1.9), Inches(0.4),
                     label, font_size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(s, x, Inches(3.5), Inches(1.9), Inches(0.8),
                     desc, font_size=10, color=SLATE, align=PP_ALIGN.CENTER)
        x += Inches(2.1)
    slide_footer(s)

    # Slide 70 — Cost Comparison
    s = prs.slides.add_slide(BLANK)
    table_slide(s, "Cost Comparison — Elenza PMS vs Alternatives",
        ["Factor", "Elenza PMS", "Generic ERP", "Custom SaaS"],
        [
            ["Upfront Cost", "One-time development", "High licensing fees", "Monthly recurring"],
            ["Monthly Cost", "Hosting only (~$5-15)", "$50-500+/month", "$100-1000+/month"],
            ["Customization", "Full source ownership", "Limited customization", "Vendor-dependent"],
            ["Data Ownership", "Your server, your data", "Vendor-controlled", "Cloud-locked"],
            ["Internet Required", "No (local network)", "Usually yes", "Always yes"],
            ["IT Team Needed", "Minimal", "Dedicated team", "Vendor support"],
            ["Production Focus", "Built for factory floor", "General purpose", "Generic workflows"],
            ["Factory-Floor UI", "Yes — table-first", "Dashboard-heavy", "Varies"],
            ["Source Code", "100% yours", "No access", "No access"],
            ["Long-Term TCO", "Lowest", "Highest", "Medium-High"],
        ],
        [Inches(2.5), Inches(3.3), Inches(3.3), Inches(3.3)])

    # Slide 71 — Competitive Advantage
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Competitive Advantage",
                 "Why Modular Interior Factories Choose Elenza PMS")
    advantages = [
        "\u2714  Purpose-built for modular interiors — not a generic tool",
        "\u2714  Machine-wise tracking with station-level login — unique in the market",
        "\u2714  Partial completion logic handles real factory scenarios",
        "\u2714  Rejection routing with mandatory reason — traceable corrections",
        "\u2714  Packing and Dispatch as separate operational portals",
        "\u2714  Planner workspace with priority, machine grouping, and Excel export",
        "\u2714  One current stage — no confusion about order status",
        "\u2714  Admin manages everything without developer help",
        "\u2714  No monthly fees — own the system forever",
        "\u2714  Runs on basic Windows hosting — no expensive infrastructure",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                advantages, font_size=14, color=SLATE)
    slide_footer(s)

    # Slide 72 — Business Benefits
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Business Benefits Summary",
                 "What You Get with Elenza PMS")
    benefits = [
        ("Faster\nOperations", "60%+ faster data entry\nwith searchable dropdowns"),
        ("Full\nVisibility", "Know exactly where\nevery order is"),
        ("Better\nDecisions", "Reports with search,\nfilter, and Excel export"),
        ("Reduced\nErrors", "Role-based access\nprevents unauthorized changes"),
        ("Complete\nAudit", "Every action logged\nwith user and timestamp"),
        ("Scalable\nGrowth", "Add machines, users,\nstations without developer"),
    ]
    x = Inches(0.4)
    for title, desc in benefits:
        add_box(s, x, Inches(1.6), Inches(1.9), Inches(1.3), title, BLUE, WHITE, 13, True)
        add_text_box(s, x, Inches(3.1), Inches(1.9), Inches(1.0),
                     desc, font_size=10, color=SLATE, align=PP_ALIGN.CENTER)
        x += Inches(2.1)
    slide_footer(s)


def build_planner_deep_dive(prs):
    # Slide 73 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "09", "Production Flow Deep Dive", "\u2699")

    # Slide 74 — Complete Order Flow
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Complete Order Status Flow",
                 "From Quotation to Dispatched")
    flow_steps = [
        "Quotation\nCreated",
        "Confirmed",
        "Optimised",
        "Procurement\nStarted",
        "Material\nReceived",
        "Production\nStarted",
        "Hot Press",
        "Cutting",
        "Edgebanding",
    ]
    flow_slide(s, "", flow_steps)
    # Second row
    flow_steps2 = ["Drilling", "QC", "Packing", "Dispatch", "Fully\nCompleted"]
    n2 = len(flow_steps2)
    box_w = Inches(1.6)
    gap = Inches(0.25)
    arrow_w = Inches(0.4)
    total = n2 * box_w + (n2 - 1) * (gap + arrow_w)
    start_x = (SLIDE_W - total) // 2
    y_box = Inches(5.0)
    y_arrow = y_box + box_w // 2 - Inches(0.15)
    for i, step in enumerate(flow_steps2):
        x = start_x + i * (box_w + gap + arrow_w)
        c = BLUE if i % 2 == 0 else BLUE_DARK
        add_box(s, x, y_box, box_w, Inches(0.9), step, c, WHITE, 12, True)
        if i < n2 - 1:
            add_arrow_right(s, x + box_w + gap // 2, y_arrow, arrow_w,
                            Inches(0.2), BLUE_DARK)
    slide_footer(s)

    # Slide 75 — Production Statuses
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Production Status Values",
                 "All Possible Order States")
    statuses = [
        ("Pending", GRAY, "Awaiting action\nat current station"),
        ("In Progress", BLUE, "Currently being\nworked on"),
        ("Completed", GREEN, "Done, moved\nto next station"),
        ("Partial Completed", ORANGE, "Partially done,\n2 stations visible"),
        ("Rejected", RED, "Sent backward\nfor correction"),
        ("Production Completed", NAVY, "All stations done,\nready for packing"),
        ("Pending Dispatch", BLUE_DARK, "Packed, awaiting\ndispatch action"),
        ("Partially Dispatched", ORANGE, "Some items\ndispatched"),
        ("Dispatched", GREEN, "Fully shipped\nand closed"),
        ("Hold", RED, "Order held,\nremarks required"),
    ]
    x = Inches(0.3)
    for i, (label, color, desc) in enumerate(statuses):
        col = i % 5
        row = i // 5
        xx = Inches(0.3) + col * Inches(2.5)
        yy = Inches(1.5) + row * Inches(2.6)
        add_box(s, xx, yy, Inches(2.2), Inches(0.8), label, color, WHITE, 11, True)
        add_text_box(s, xx, yy + Inches(0.9), Inches(2.2), Inches(0.8),
                     desc, font_size=10, color=SLATE, align=PP_ALIGN.CENTER)
    slide_footer(s)

    # Slide 76 — Production Sequence Configuration
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Production Sequence Configuration",
                 "Admin Manages Without Developer")
    capabilities = [
        ("\u2795", "Add New Station", "Insert new machine or\nprocess step into sequence"),
        ("\u270F", "Rename Station", "Update station names\nto match factory floor"),
        ("\u2716", "Disable Station", "Remove from active\nsequence while preserving data"),
        ("\u21C5", "Reorder", "Drag-and-drop or\nnumeric reordering"),
        ("\u279C", "Insert Between", "Place new station\nbetween existing ones"),
    ]
    icon_card_row(s, capabilities, y= Inches(1.5))
    add_text_box(s, Inches(0.5), Inches(6.0), Inches(12), Inches(0.6),
                 "\u26A0  Sequence changes do NOT break old records — historical data preserved",
                 font_size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    slide_footer(s)

    # Slide 77 — Planner Movement Actions
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Planner Movement Actions",
                 "Assign, Resequence, Reapprove")
    actions = [
        ("\u25B6", "Assign to Station", "Move order to specific\nmachine in sequence"),
        ("\u21C5", "Resequence", "Change station order\nfor specific orders"),
        ("\u2714", "Reapprove", "After rejection, planner\nreintroduces to production"),
        ("\u2611", "Set Priority", "High, Medium, Low,\nor blank"),
        ("\u231B", "View History", "Full lifecycle from\nplanner context"),
    ]
    icon_card_row(s, actions, y= Inches(1.5))
    slide_footer(s)


def build_deployment(prs):
    # Slide 78 — Deployment Architecture
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Deployment Architecture",
                 "Simple, Reliable, Maintainable")
    add_rect(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.2), LIGHT_BG, BORDER, 0.02)
    layers = [
        ("IIS / ASP.NET 4.8 Hosting", Inches(1.8), BLUE),
        ("Windows Server / Shared Hosting", Inches(2.8), BLUE_DARK),
        ("Microsoft Access Database (App_Data)", Inches(3.8), NAVY),
        ("FTP Deployment (Passive Binary)", Inches(4.8), GRAY),
        ("Backup & Rollback System", Inches(5.8), BLUE),
    ]
    for label, y, color in layers:
        add_box(s, Inches(1), y, Inches(11.3), Inches(0.8), label, color, WHITE, 14, True)
    slide_footer(s)

    # Slide 79 — Deployment Process
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Deployment Process",
                 "Verified Upload Workflow")
    steps = [
        "Read Back\nLive File",
        "Create\nLocal Backup",
        "Patch\nVerified Copy",
        "Upload\nvia FTP",
        "Compare\nSHA-256",
        "Verify\nHTTP 200",
    ]
    flow_slide(s, "", steps)
    add_text_box(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.8),
                 "\u2714  Every live file backed up before modification\n"
                 "\u2714  SHA-256 hash verified after upload\n"
                 "\u2714  Public endpoint tested for HTTP 200 response",
                 font_size=12, color=SLATE, align=PP_ALIGN.CENTER)
    slide_footer(s)

    # Slide 80 — Rollback Capability
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Rollback Capability",
                 "Snapshot System for Safe Live Changes")
    features = [
        ("\u2611", "Snapshot", "Create timestamped backup\nbefore any live change"),
        ("\u2714", "Revert", "Restore any snapshot\nwith one command"),
        ("\u231B", "Deploy Log", "Full history of uploads\nand reverts"),
        ("\u2611", "Offline Mode", "App goes briefly offline\nduring revert"),
        ("\u2611", "Auto Cleanup", "Keep two newest\nrelevant backup sets"),
    ]
    icon_card_row(s, features, y= Inches(1.5))
    slide_footer(s)


def build_implementation(prs):
    # Slide 81 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "10", "Implementation Plan", "\u279C")

    # Slide 82 — Implementation Phases
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Implementation Phases",
                 "Structured Rollout for Factory Operations")
    phases = [
        ("Phase 1", "Core System", "Login, Dealers, Quotations,\nOrders, Basic Production", GREEN),
        ("Phase 2", "Production Flow", "Machine Tracking, Packing,\nDispatch, Sequence Config", BLUE),
        ("Phase 3", "Planner", "Planner Workspace, Priority,\nMachine Wise, Excel Export", BLUE_DARK),
        ("Phase 4", "Reports", "All Reports, Filters, Export,\nLifecycle, Audit Trail", NAVY),
    ]
    x = Inches(0.4)
    for label, title, desc, color in phases:
        add_box(s, x, Inches(1.6), Inches(3.0), Inches(0.7), label, color, WHITE, 14, True)
        add_box(s, x, Inches(2.4), Inches(3.0), Inches(0.6), title, LIGHT_BG, NAVY, 14, True)
        add_text_box(s, x, Inches(3.1), Inches(3.0), Inches(1.0),
                     desc, font_size=11, color=SLATE, align=PP_ALIGN.CENTER)
        x += Inches(3.2)
    add_text_box(s, Inches(0.5), Inches(4.5), Inches(12), Inches(2.0),
                 "\u2714  Each phase is independently deployable\n"
                 "\u2714  Phase 1+2 covers daily factory operations\n"
                 "\u2714  Phase 3 adds planning intelligence\n"
                 "\u2714  Phase 4 completes reporting and audit",
                 font_size=14, color=NAVY)
    slide_footer(s)

    # Slide 83 — Timeline
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Estimated Timeline",
                 "Factory Operational in Weeks, Not Months")
    milestones = [
        ("Week 1-2", "Core Setup", "Login, dealers,\nquotations, orders"),
        ("Week 3-4", "Production", "Machine tracking,\nsequence, basic flow"),
        ("Week 5-6", "Packing & Dispatch", "Packing portal,\ndispatch workflow"),
        ("Week 7-8", "Planner", "Planner workspace,\npriority, exports"),
        ("Week 9-10", "Reports & Audit", "All reports,\nlifecycle, audit trail"),
        ("Week 11-12", "Training & Go-Live", "User training,\ndata migration, go-live"),
    ]
    x = Inches(0.3)
    for week, title, desc in milestones:
        add_box(s, x, Inches(1.6), Inches(1.9), Inches(0.6), week, BLUE, WHITE, 11, True)
        add_text_box(s, x, Inches(2.3), Inches(1.9), Inches(0.4),
                     title, font_size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(s, x, Inches(2.7), Inches(1.9), Inches(0.8),
                     desc, font_size=10, color=SLATE, align=PP_ALIGN.CENTER)
        x += Inches(2.1)
    slide_footer(s)

    # Slide 84 — Data Migration
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Data Migration Approach",
                 "Existing Data Preserved and Imported")
    steps = [
        "Export\nExisting Data",
        "Map to\nNew Schema",
        "Import\nDealers",
        "Import\nOrders",
        "Import\nUsers",
        "Verify\nAll Records",
    ]
    flow_slide(s, "", steps)
    add_text_box(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.8),
                 "\u2714  Existing dealer, order, and user data can be imported\n"
                 "\u2714  Master values (customer type, order type) pre-configured\n"
                 "\u2714  Historical records preserved in new system",
                 font_size=12, color=SLATE, align=PP_ALIGN.CENTER)
    slide_footer(s)

    # Slide 85 — Go-Live Checklist
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Go-Live Checklist",
                 "Prepared for Production Deployment")
    checklist = [
        "\u2714  All modules tested with real factory scenarios",
        "\u2714  User accounts created with correct roles",
        "\u2714  Machine stations configured in correct sequence",
        "\u2714  Customer types and order types loaded",
        "\u2714  Existing dealer data imported and verified",
        "\u2714  Production planner configured with active orders",
        "\u2714  Reports tested with sample data",
        "\u2714  Audit trail verified for all operations",
        "\u2714  User training completed for all roles",
        "\u2714  Backup system tested and verified",
        "\u2714  Rollback procedure documented",
        "\u2714  Admin trained on master management",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                checklist, font_size=14, color=SLATE, bullet_char="")
    slide_footer(s)

    # Slide 86 — Support Model
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Support Model",
                 "Post-Deployment Assistance")
    support = [
        ("\u260E", "Technical Support", "Bug fixes and\nissue resolution"),
        ("\u2699", "Configuration", "Master setup,\nsequence changes"),
        ("\u21E9", "Updates", "Feature additions\nas needed"),
        ("\u2611", "Training", "New user training\nand documentation"),
        ("\u231B", "Monitoring", "Performance and\nusage monitoring"),
    ]
    icon_card_row(s, support, y= Inches(1.5))
    slide_footer(s)


def build_roles_deep_dive(prs):
    # Slide 87 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "11", "User Roles Deep Dive", "\u2699")

    # Slide 88 — Admin Capabilities
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Admin / Super User",
                 "Complete System Control")
    capabilities = [
        "Manage all users (create, edit, activate/deactivate)",
        "Manage all dealers (add, edit, change customer type)",
        "Manage all dropdown masters (customer type, order type)",
        "Manage machine / station master and production sequence",
        "Manage vendor master for procurement",
        "View and edit all orders at any stage",
        "View all reports across all modules",
        "Correct wrongly updated statuses",
        "Access full audit trail for any order",
        "Override machine station assignments",
        "Change production sequence without developer",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                capabilities, font_size=14, color=SLATE, bullet_char="\u2714")
    slide_footer(s)

    # Slide 89 — Data Entry User
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Data Entry User",
                 "Front-Line Order Entry")
    can = [
        "\u2714  Add and manage dealer data",
        "\u2714  Create quotation entries with all fields",
        "\u2714  Confirm orders",
        "\u2714  Add optimisation details",
        "\u2714  Add procurement details",
        "\u2714  View relevant reports",
    ]
    cannot = [
        "\u2716  Cannot manage users",
        "\u2716  Cannot change machine sequence",
        "\u2716  Cannot access planner workspace",
        "\u2716  Cannot modify master settings",
        "\u2716  Cannot update machine station orders",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.0),
                can, font_size=14, color=GREEN, bullet_char="")
    add_bullets(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.0),
                cannot, font_size=14, color=RED, bullet_char="")
    slide_footer(s)

    # Slide 90 — Machine User
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Machine User",
                 "Station-Specific Production Queue")
    can = [
        "\u2714  See only orders assigned to their station",
        "\u2714  Mark orders as Completed",
        "\u2714  Mark orders as Partial Completed",
        "\u2714  Mark orders as Rejected",
        "\u2714  Add remarks for each action",
        "\u2714  View date and time of last update",
    ]
    cannot = [
        "\u2716  Cannot edit dealer data",
        "\u2716  Cannot edit quotation details",
        "\u2716  Cannot modify procurement data",
        "\u2716  Cannot access other station orders",
        "\u2716  Cannot change production sequence",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(5.0),
                can, font_size=14, color=GREEN, bullet_char="")
    add_bullets(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(5.0),
                cannot, font_size=14, color=RED, bullet_char="")
    slide_footer(s)

    # Slide 91 — Planner User
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Production Planner User",
                 "Planning Intelligence & Priority Management")
    capabilities = [
        "\u2714  Open separate planner workspace page after login",
        "\u2714  View all active planning orders (Confirmed to Packing)",
        "\u2714  View machine-wise grouped planning tables",
        "\u2714  Assign priority: High, Medium, Low, or blank",
        "\u2714  View current stage for every active order",
        "\u2714  Export each planner table to Excel",
        "\u2714  Export machine-wise data sorted by confirmation date descending",
        "\u2714  Open order history from planner context",
        "\u2714  Cannot treat one order as having more than one current stage",
        "\u2714  Cannot access data entry or master editing",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                capabilities, font_size=14, color=SLATE, bullet_char="")
    slide_footer(s)


def build_database(prs):
    # Slide 92 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "12", "Database Design", "\u2630")

    # Slide 93 — Database Tables
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Microsoft Access Database Schema",
                 "19 Core Tables — Complete Data Model")
    tables = [
        ["tbl_users", "User accounts, roles, and station assignments"],
        ["tbl_roles", "Role definitions and permissions"],
        ["tbl_dealers", "Dealer profiles and contact details"],
        ["tbl_quotations", "Quotation and enquiry records"],
        ["tbl_orders", "Confirmed orders with customer and type data"],
        ["tbl_order_confirmations", "Order confirmation events and dates"],
        ["tbl_optimisation", "Board optimisation and raw material details"],
        ["tbl_procurement", "Purchase order and material receipt records"],
        ["tbl_procurement_items", "Individual items against procurement"],
        ["tbl_vendors", "Vendor master for procurement"],
        ["tbl_order_types", "Order type master values"],
        ["tbl_customer_types", "Customer type master values"],
        ["tbl_machines", "Machine and station master"],
        ["tbl_machine_sequence", "Production sequence ordering"],
        ["tbl_production_tracking", "Current production status per order"],
        ["tbl_production_history", "Historical production movements"],
        ["tbl_dispatch", "Dispatch details and status"],
        ["tbl_audit_logs", "Complete audit trail for all actions"],
        ["tbl_dropdown_masters", "Configurable dropdown values"],
    ]
    table_slide(s, "", ["Table", "Purpose"], tables,
                [Inches(3.5), Inches(8.5)])
    slide_footer(s)

    # Slide 94 — Database Relationships
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Database Relationships",
                 "Key Table Connections")
    relationships = [
        "\u2022  tbl_quotations.dealer_id  ->  tbl_dealers.dealer_id",
        "\u2022  tbl_orders.quotation_id  ->  tbl_quotations.quotation_id",
        "\u2022  tbl_order_confirmations.order_id  ->  tbl_orders.order_id",
        "\u2022  tbl_optimisation.order_id  ->  tbl_orders.order_id",
        "\u2022  tbl_procurement.order_id  ->  tbl_orders.order_id",
        "\u2022  tbl_procurement_items.procurement_id  ->  tbl_procurement.procurement_id",
        "\u2022  tbl_production_tracking.order_id  ->  tbl_orders.order_id",
        "\u2022  tbl_production_tracking.station_id  ->  tbl_machines.machine_id",
        "\u2022  tbl_production_history.order_id  ->  tbl_orders.order_id",
        "\u2022  tbl_dispatch.order_id  ->  tbl_orders.order_id",
        "\u2022  tbl_audit_logs.order_id  ->  tbl_orders.order_id (when applicable)",
        "\u2022  tbl_machine_sequence.machine_id  ->  tbl_machines.machine_id",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                relationships, font_size=13, color=SLATE, bullet_char="")
    slide_footer(s)


def build_future_roadmap(prs):
    # Slide 95 — Section
    s = prs.slides.add_slide(BLANK)
    section_title_slide(s, "13", "Future Roadmap", "\u2606")

    # Slide 96 — Planned Features
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Planned Enhancements",
                 "Continuous Improvement Based on Factory Needs")
    features = [
        ("\u2611", "Production Batches", "Group multiple orders\nfor batch processing"),
        ("\u2611", "QR Labels", "Print QR labels for\norder and batch tracking"),
        ("\u2611", "Barcode Scanning", "Scan QR/barcodes at\nmachine stations"),
        ("\u2611", "Inventory Valuation", "Track material costs\nand inventory value"),
        ("\u2611", "Accounting Integration", "Connect with existing\naccounting systems"),
        ("\u2611", "Mobile App", "Native mobile for\nmachine operators"),
        ("\u2611", "WhatsApp Alerts", "Order status updates\nvia WhatsApp"),
        ("\u2611", "Customer Portal", "Dealer self-service\norder tracking"),
        ("\u2611", "Advanced Analytics", "Business intelligence\ndashboards"),
    ]
    icon_card_row(s, features, y= Inches(1.5))
    slide_footer(s)

    # Slide 97 — Roadmap Timeline
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Roadmap Timeline",
                 "Phased Enhancement Delivery")
    milestones = [
        ("Current", "Core PMS operational\nwith planner, packing,\ndispatch", GREEN),
        ("Q3 2026", "Source consolidation,\nautomated tests,\nsecurity hardening", BLUE),
        ("Q4 2026", "Production batches,\nQR labels,\nbarcode scanning", BLUE_DARK),
        ("2027", "Mobile app,\nWhatsApp integration,\ncustomer portal", NAVY),
    ]
    x = Inches(0.4)
    for period, desc, color in milestones:
        add_box(s, x, Inches(1.8), Inches(2.8), Inches(0.8), period, color, WHITE, 16, True)
        add_text_box(s, x, Inches(2.8), Inches(2.8), Inches(1.2),
                     desc, font_size=12, color=SLATE, align=PP_ALIGN.CENTER)
        x += Inches(3.1)
    slide_footer(s)


def build_closing(prs):
    # Slide 98 — Summary
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Summary",
                 "Elenza PMS — Complete Production Management")
    summary = [
        "\u2714  13 core modules covering quotation to dispatch",
        "\u2714  Machine-wise production tracking with station login",
        "\u2714  Partial completion and rejection handling",
        "\u2714  Dedicated Packing and Dispatch portals",
        "\u2714  Production Planner with priority and machine grouping",
        "\u2714  12+ reports with search, filter, and Excel export",
        "\u2714  Complete audit trail for every action",
        "\u2714  Admin-managed masters, sequences, and users",
        "\u2714  Role-based access enforced server-side",
        "\u2714  Simple, proven technology stack",
        "\u2714  One-time investment — no monthly fees",
        "\u2714  Built specifically for modular interior B2B operations",
    ]
    add_bullets(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5),
                summary, font_size=14, color=SLATE, bullet_char="")
    slide_footer(s)

    # Slide 99 — Next Steps
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "Next Steps",
                 "How to Get Started")
    steps = [
        ("1", "Discovery Call", "Discuss your factory\noperations and needs"),
        ("2", "Requirements", "Map your specific\nworkflows to modules"),
        ("3", "Configuration", "Set up users, machines,\nsequences, masters"),
        ("4", "Data Migration", "Import existing dealers,\norders, and users"),
        ("5", "Training", "Role-based training\nfor all user types"),
        ("6", "Go Live", "Deploy and start\nusing the system"),
    ]
    x = Inches(0.4)
    for num, title, desc in steps:
        add_circle(s, x + Inches(0.65), Inches(1.5), Inches(0.9), BLUE, num, WHITE, 24)
        add_text_box(s, x, Inches(2.6), Inches(2.0), Inches(0.4),
                     title, font_size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(s, x, Inches(3.0), Inches(2.0), Inches(0.8),
                     desc, font_size=11, color=SLATE, align=PP_ALIGN.CENTER)
        x += Inches(2.1)
    add_text_box(s, Inches(0.5), Inches(4.5), Inches(12), Inches(2.0),
                 "\n\n"
                 "Let's discuss how Elenza PMS can transform\n"
                 "your modular interior factory operations.",
                 font_size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    slide_footer(s)

    # Slide 100 — Thank You / Contact
    s = prs.slides.add_slide(BLANK)
    add_bg(s, NAVY)
    add_rect(s, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5), BLUE, None, 0.02)
    add_text_box(s, Inches(2), Inches(2.0), Inches(9.3), Inches(1.0),
                 "Thank You", font_size=48, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER)
    add_rect(s, Inches(5), Inches(3.1), Inches(3.3), Inches(0.04), WHITE)
    add_text_box(s, Inches(2), Inches(3.4), Inches(9.3), Inches(0.6),
                 "ElenzaIndia.com Production Management System",
                 font_size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(2), Inches(4.2), Inches(9.3), Inches(0.5),
                 "Questions?  Let's Discuss.",
                 font_size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(2), Inches(5.0), Inches(9.3), Inches(0.4),
                 "www.elenzaindia.com",
                 font_size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(2), Inches(6.2), Inches(9.3), Inches(0.4),
                 "Built with Purpose  |  Designed for Factory Floors",
                 font_size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# MAIN — Build all slides
# ═══════════════════════════════════════════════════════════════════

def main():
    print("Building Elenza PMS Pitch Deck — 100 slides...")
    print("  Section 1:  Cover & Agenda          (slides 1-3)")
    build_cover(prs)
    print("  Section 2:  Problem Statement       (slides 4-9)")
    build_problem(prs)
    print("  Section 3:  Solution Overview        (slides 10-13)")
    build_solution(prs)
    print("  Section 4:  Modules Overview         (slides 14-16)")
    build_modules_overview(prs)
    print("  Section 5:  Login & Roles            (slides 17-19)")
    build_login_module(prs)
    print("  Section 6:  Dealer Management        (slides 20-21)")
    build_dealer_module(prs)
    print("  Section 7:  Quotation                (slides 22-23)")
    build_quotation_module(prs)
    print("  Section 8:  Order Confirmation       (slides 24-25)")
    build_confirmation_module(prs)
    print("  Section 9:  Optimisation             (slides 26-27)")
    build_optimisation_module(prs)
    print("  Section 10: Procurement              (slides 28-29)")
    build_procurement_module(prs)
    print("  Section 11: Production Tracking      (slides 30-36)")
    build_production_tracking(prs)
    print("  Section 12: Packing                  (slides 37-39)")
    build_packing(prs)
    print("  Section 13: Dispatch                 (slides 40-41)")
    build_dispatch(prs)
    print("  Section 14: Production Planner       (slides 42-48)")
    build_planner(prs)
    print("  Section 15: Reports                  (slides 49-53)")
    build_reports(prs)
    print("  Section 16: Masters                  (slides 54-60)")
    build_masters(prs)
    print("  Section 17: Audit Trail              (slides 61-62)")
    build_audit(prs)
    print("  Section 18: Technology               (slides 63-67)")
    build_technology(prs)
    print("  Section 19: Value Proposition        (slides 68-72)")
    build_value_proposition(prs)
    print("  Section 20: Production Flow Deep     (slides 73-77)")
    build_planner_deep_dive(prs)
    print("  Section 21: Deployment               (slides 78-80)")
    build_deployment(prs)
    print("  Section 22: Implementation           (slides 81-86)")
    build_implementation(prs)
    print("  Section 23: Roles Deep Dive          (slides 87-91)")
    build_roles_deep_dive(prs)
    print("  Section 24: Database                 (slides 92-94)")
    build_database(prs)
    print("  Section 25: Future Roadmap           (slides 95-97)")
    build_future_roadmap(prs)
    print("  Section 26: Closing                  (slides 98-100)")
    build_closing(prs)

    total = len(prs.slides)
    print(f"\nTotal slides generated: {total}")

    out = "ElenzaPMS_PitchDeck.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print("Done!")


if __name__ == "__main__":
    main()
